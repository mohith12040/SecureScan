import sys
import os
import subprocess

# Self-installation check for python-pptx
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing dynamically...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
    except Exception as e:
        print(f"Error installing python-pptx: {e}")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -------------------------------------------------------------
# Color Palette Definition (Premium Dark Cyber)
# -------------------------------------------------------------
COLOR_BG = RGBColor(10, 15, 29)          # Deep Charcoal/Matte Black (#0A0F1D)
COLOR_TEXT_PRIMARY = RGBColor(255, 255, 255) # Pure White
COLOR_TEXT_SECONDARY = RGBColor(178, 190, 195) # Cool Gray (#B2BEC3)
COLOR_CYAN = RGBColor(0, 242, 254)       # Cyber Cyan (#00F2FE)
COLOR_TEAL = RGBColor(9, 132, 227)       # Deep Teal (#0984E3)
COLOR_MUTED_BLUE = RGBColor(25, 35, 60)  # Lighter Panel Blue (#19233C)
COLOR_WARNING = RGBColor(255, 118, 117)   # Muted red/coral (#FF7675)

def apply_background(slide):
    """Applies a solid deep-dark background to the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title, category="SECURESCAN"):
    """Adds a standard header layout with category tracker, title, and divider line."""
    # Category Tracker Text
    tracker_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_tr = tracker_box.text_frame
    tf_tr.word_wrap = True
    tf_tr.margin_left = tf_tr.margin_top = tf_tr.margin_right = tf_tr.margin_bottom = 0
    p_tr = tf_tr.paragraphs[0]
    p_tr.text = category.upper()
    p_tr.font.name = "Arial"
    p_tr.font.size = Pt(10)
    p_tr.font.bold = True
    p_tr.font.color.rgb = COLOR_TEAL
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    # Dynamic divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEAL
    line.line.color.rgb = COLOR_TEAL

def create_unified_card(slide, left, top, width, height, title="", bullets=None, title_color=COLOR_CYAN, fill_color=COLOR_MUTED_BLUE, font_size=13, space_after=5):
    """Draws a single card panel using standard RECTANGLE and renders both title and bullet points inside its own text_frame.
    This guarantees that the text and card background/border are unified in a single shape and can never drift.
    """
    # Use standard RECTANGLE to prevent ROUNDED_RECTANGLE warp/drift bugs in third-party viewers
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLOR_TEAL
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.25)
    
    first_paragraph = True
    
    # Render title if provided
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(10)
        first_paragraph = False
        
    # Render bullet points
    if bullets:
        for item in bullets:
            lines = item.split('\n')
            for line in lines:
                line_str = line.strip()
                if not line_str:
                    continue
                
                is_sub = False
                if line_str.startswith(('-', '*', '•')):
                    is_sub = True
                    line_str = line_str[1:].strip()
                    
                if first_paragraph:
                    p = tf.paragraphs[0]
                    first_paragraph = False
                else:
                    p = tf.add_paragraph()
                    
                p.font.name = "Arial"
                
                if is_sub:
                    # Tab separation + level 1 indent gives perfect hanging indents
                    p.text = "–\t" + line_str
                    p.font.size = Pt(font_size - 1)
                    p.font.color.rgb = COLOR_TEXT_SECONDARY
                    p.space_after = Pt(space_after - 1)
                    p.level = 1
                else:
                    p.text = "•\t" + line_str
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = COLOR_TEXT_PRIMARY
                    p.space_after = Pt(space_after)
                    p.level = 0
                    
    return shape

def create_stage_card(slide, left, top, width, height, num, name, desc_lines, fill_color=COLOR_MUTED_BLUE):
    """Creates a stage block inside a single rectangle shape with a text divider, ensuring elements never drift."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLOR_TEAL
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.2)
    
    # Stage Number
    p_num = tf.paragraphs[0]
    p_num.text = num
    p_num.font.name = "Arial"
    p_num.font.size = Pt(11)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_TEAL
    p_num.alignment = PP_ALIGN.CENTER
    p_num.space_after = Pt(4)
    
    # Stage Name
    p_name = tf.add_paragraph()
    p_name.text = name
    p_name.font.name = "Arial"
    p_name.font.size = Pt(18)
    p_name.font.bold = True
    p_name.font.color.rgb = COLOR_CYAN
    p_name.alignment = PP_ALIGN.CENTER
    p_name.space_after = Pt(8)
    
    # Text-based divider line
    p_div = tf.add_paragraph()
    p_div.text = "────────"
    p_div.font.name = "Arial"
    p_div.font.size = Pt(10)
    p_div.font.color.rgb = COLOR_TEAL
    p_div.alignment = PP_ALIGN.CENTER
    p_div.space_after = Pt(10)
    
    # Description Lines
    for line in desc_lines:
        p_desc = tf.add_paragraph()
        p_desc.text = line.strip()
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_PRIMARY
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.space_after = Pt(4)
        
    return shape

def create_metric_card(slide, left, top, width, height, metric_text, label_text, desc_text, fill_color=COLOR_MUTED_BLUE):
    """Draws a metric score display box unified in a single shape to avoid offset drift."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLOR_TEAL
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.4)
    tf.margin_bottom = Inches(0.2)
    
    # Score Metric
    p_met = tf.paragraphs[0]
    p_met.text = metric_text
    p_met.font.name = "Arial"
    p_met.font.size = Pt(72)
    p_met.font.bold = True
    p_met.font.color.rgb = COLOR_CYAN
    p_met.alignment = PP_ALIGN.CENTER
    p_met.space_after = Pt(10)
    
    # Metric Label
    p_lab = tf.add_paragraph()
    p_lab.text = label_text
    p_lab.font.name = "Arial"
    p_lab.font.size = Pt(16)
    p_lab.font.bold = True
    p_lab.font.color.rgb = COLOR_TEXT_PRIMARY
    p_lab.alignment = PP_ALIGN.CENTER
    p_lab.space_after = Pt(15)
    
    # Detail Text
    p_desc = tf.add_paragraph()
    p_desc.text = desc_text
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_TEXT_SECONDARY
    p_desc.alignment = PP_ALIGN.CENTER
    p_desc.space_after = Pt(2)
    
    return shape

def create_roadmap_card(slide, left, top, width, height, step, name, desc, fill_color=COLOR_MUTED_BLUE):
    """Draws a unified roadmap panel inside a single rectangle shape with a text divider."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLOR_TEAL
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.2)
    
    # Step header
    p_step = tf.paragraphs[0]
    p_step.text = step
    p_step.font.name = "Arial"
    p_step.font.size = Pt(11)
    p_step.font.bold = True
    p_step.font.color.rgb = COLOR_CYAN
    p_step.space_after = Pt(4)
    
    # Step Title
    p_name = tf.add_paragraph()
    p_name.text = name
    p_name.font.name = "Arial"
    p_name.font.size = Pt(18)
    p_name.font.bold = True
    p_name.font.color.rgb = COLOR_TEXT_PRIMARY
    p_name.space_after = Pt(10)
    
    # Text-based divider line
    p_div = tf.add_paragraph()
    p_div.text = "────────────"
    p_div.font.name = "Arial"
    p_div.font.size = Pt(10)
    p_div.font.color.rgb = COLOR_TEAL
    p_div.space_after = Pt(12)
    
    # Description
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_TEXT_SECONDARY
    p_desc.space_after = Pt(6)
    
    return shape

def main():
    prs = Presentation()
    # Configure 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Cover)
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    
    # Giant Cyan Decorative Accent Line on the Left
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.12), Inches(3.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_CYAN
    accent_bar.line.color.rgb = COLOR_CYAN
    
    # Title Text Frame
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "SECURESCAN"
    p.font.name = "Arial"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Driven Vulnerability Management Platform"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_TEXT_PRIMARY
    p2.space_after = Pt(8)
    
    p3 = tf.add_paragraph()
    p3.text = "Merging Network Probing with Random Forest Classifiers for Enterprise Audits"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_TEXT_SECONDARY
    
    # Sub-footer details
    footer_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(8.0), Inches(0.8))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
    
    lines = ["Presenter: Cyber Engineering Team", "Technical Portfolio Showcase"]
    for idx, line in enumerate(lines):
        if idx == 0:
            p_f = tf_f.paragraphs[0]
        else:
            p_f = tf_f.add_paragraph()
        p_f.text = line
        p_f.font.name = "Arial"
        p_f.font.size = Pt(12)
        p_f.font.color.rgb = COLOR_TEAL
        p_f.space_after = Pt(2)

    # -------------------------------------------------------------
    # SLIDE 2: The Cybersecurity Challenge
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "The Cybersecurity Challenge")
    
    # Descriptive Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Modern enterprises face critical bottlenecks in vulnerability discovery and analysis:"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_TEXT_SECONDARY
    
    # Left Grid (Problem 1 & 2)
    bullets_p1 = [
        "Analyzing open system ports and services relies heavily on senior security specialists.",
        "Manual calculations of CVSS scores lead to significant assessment latency."
    ]
    create_unified_card(slide, Inches(0.8), Inches(2.2), Inches(5.6), Inches(2.1), "1. Excessive Manual Overhead", bullets_p1, COLOR_WARNING)
    
    bullets_p2 = [
        "Traditional tools dump raw, unstructured text files filled with complex version headers.",
        "System administrators struggle to extract immediate actionable intelligence."
    ]
    create_unified_card(slide, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.1), "2. Raw Scan Output Complexity", bullets_p2, COLOR_WARNING)
    
    # Right Grid (Problem 3 & 4)
    bullets_p3 = [
        "Threat classification uses rigid, static rule matrices.",
        "Fails to incorporate contextual feature weighting like exposed scope or overall port counts."
    ]
    create_unified_card(slide, Inches(6.8), Inches(2.2), Inches(5.7), Inches(2.1), "3. Absence of Intelligent Context", bullets_p3, COLOR_WARNING)
    
    bullets_p4 = [
        "Scanners outline problems but fail to offer instant, targeted configuration rules.",
        "Creates a severe delay between vulnerability detection and defense execution."
    ]
    create_unified_card(slide, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.1), "4. Disconnected Defensive Action", bullets_p4, COLOR_WARNING)

    # -------------------------------------------------------------
    # SLIDE 3: The SecureScan Solution
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "The SecureScan Solution")
    
    # Top subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "An automated ecosystem that unifies live network probing, AI diagnostics, and rapid reporting:"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_TEXT_SECONDARY

    # 4 Quadrants for the Solution
    cols = [
        {"l": Inches(0.8), "t": Inches(2.2), "w": Inches(5.6), "h": Inches(2.1), "title": "Automated Network Probing", "bullets": ["Nmap deep probing core to scan services, active ports, and OS versions.", "Robust simulated scanner fallback for offline sandboxed audits."]},
        {"l": Inches(6.8), "t": Inches(2.2), "w": Inches(5.7), "h": Inches(2.1), "title": "Machine Learning Severity Engine", "bullets": ["Supervised Random Forest Classifier predicting danger scopes under 5ms.", "Accurately labels risk vectors (Critical, High, Medium, Low)."]},
        {"l": Inches(0.8), "t": Inches(4.5), "w": Inches(5.6), "h": Inches(2.1), "title": "Intelligent Threat Remediation", "bullets": ["Maps exposed system ports straight to explicit security action plans.", "Provides instant configuration fixes and firewall containment rules."]},
        {"l": Inches(6.8), "t": Inches(4.5), "w": Inches(5.7), "h": Inches(2.1), "title": "Executive Report Compiler", "bullets": ["Programmatic ReportLab engine generating multi-page executive PDFs.", "Includes clean telemetry tables, indicators, and prioritized checklists."]}
    ]
    
    for c in cols:
        create_unified_card(slide, c["l"], c["t"], c["w"], c["h"], c["title"], c["bullets"], COLOR_CYAN)

    # -------------------------------------------------------------
    # SLIDE 4: System Architecture & Data Flow
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "System Architecture & Data Flow")
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "How SecureScan moves seamlessly from raw target input to compiled executive audits:"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_TEXT_SECONDARY

    # 4 logical stages
    stages = [
        {"num": "STAGE 01", "name": "Scan Engine", "desc": ["User submits IP target.", "System queries Nmap or", "triggers sandbox simulator."], "x": Inches(0.8)},
        {"num": "STAGE 02", "name": "Risk Analysis", "desc": ["Calculates threat weights.", "Parses exposure scopes and", "dangerous port flags."], "x": Inches(3.8)},
        {"num": "STAGE 03", "name": "AI Predictor", "desc": ["Loads Scikit-Learn pipeline.", "Predicts risk (98.2% Accuracy)", "with confidence score."], "x": Inches(6.8)},
        {"num": "STAGE 04", "name": "Audit Synthesis", "desc": ["Generates PDF via ReportLab.", "Stores historical records", "and streams data to UI."], "x": Inches(9.8)}
    ]
    
    for s in stages:
        create_stage_card(slide, s["x"], Inches(2.3), Inches(2.7), Inches(4.3), s["num"], s["name"], s["desc"])
        
    # Draw arrow shapes between panels
    for i in range(3):
        arrow_x = Inches(3.5 + i * 3.0)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(4.2), Inches(0.3), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_CYAN
        arrow.line.fill.background()

    # -------------------------------------------------------------
    # SLIDE 5: The AI Severity Classification Engine
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "Machine Learning Severity Engine")
    
    # Left: Large Metrics Panel
    desc_left = "Trained and cross-validated on 2,500 distinct synthetic vulnerability vectors mapping port weights, outdated protocols, and risks."
    create_metric_card(slide, Inches(0.8), Inches(2.0), Inches(4.5), Inches(4.6), "98.2%", "TEST ACCURACY SCORE\n(Random Forest Classifier)", desc_left)
    
    # Right: Methodology and Features
    bullets_ml = [
        "Core Feature Inputs Analyzed:\n- Open Ports Count: Cumulative network entrypoints.\n- Critical Port Flag: Exposure of protocols like SMB (445), RDP (3389).\n- Outdated Services Count: Presence of obsolete header versions.\n- Risk Score: Contextual custom-weighted mathematical vector.",
        "Model Regularization:\n- Supervised Random Forest Classifier with balanced weights.\n- Enforced max-depth boundary rules to prevent system overfitting.",
        "Sub-Millisecond Inference speed achieved using optimized joblib serialization modules (`model.joblib`)."
    ]
    create_unified_card(slide, Inches(5.8), Inches(2.0), Inches(6.7), Inches(4.6), "Methodology & Feature Vectors", bullets_ml, COLOR_CYAN, space_after=4)

    # -------------------------------------------------------------
    # SLIDE 6: Premium Cyber-Operations Console (UI/UX)
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "Cyber-Operations Telemetry Console")
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Premium design engineered for modern security analysts:"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_TEXT_SECONDARY

    # Columns representing UI Features
    ui_cols = [
        {"title": "Glassmorphism UI", "bullets": ["Frosted translucent panels using customized CSS dark variables.", "Highly responsive layout adaptively mapping from desktop screens to tablets.", "Crisp hover indicators and smooth page-slide interactions."]},
        {"title": "Live Scanning Feed", "bullets": ["Real-time terminal diagnostic logs rendered dynamically.", "Typewriter text animations during active network audits.", "Immediate visual transition from scan completion to AI classifier."]},
        {"title": "Interactive Telemetry", "bullets": ["Integrated Chart.js engines mapping scanner history trends.", "Color-coded vulnerability ratios (Red, Orange, Blue, Green).", "Clean pagination controls managing database audit records."]}
    ]
    
    for i, col in enumerate(ui_cols):
        left_pos = Inches(0.8 + i * 4.0)
        create_unified_card(slide, left_pos, Inches(2.2), Inches(3.7), Inches(4.4), col["title"], col["bullets"], COLOR_CYAN, space_after=6)

    # -------------------------------------------------------------
    # SLIDE 7: Actionable Mitigation & PDF Auditing
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "Actionable Mitigations & PDF Auditing")
    
    # Left Side: PDF Synthesis
    pdf_bullets = [
        "ReportLab PDF Compiler:\n- Generates structured, production-grade vulnerability audit reports dynamically on target scan completion.",
        "Professional Styling:\n- Features structured header tables, colored severity badges, customized grid columns, and dynamic page numbering.",
        "Self-Contained Portability:\n- All generated PDFs are instantly downloadable and archived automatically inside safe local directory blocks."
    ]
    create_unified_card(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6), "Programmatic Executive PDF Audits", pdf_bullets, COLOR_CYAN, space_after=8)
    
    # Right Side: Remediation suggestions
    rem_bullets = [
        "Automated Mapping Pipeline:\n- System maps open target network ports directly to optimized cybersecurity remediation checklists.",
        "Concrete Configuration Rules:\n- Provides exact mitigation commands (e.g. disabling unauthenticated bindings, block rules, protocol upgrades).",
        "Prioritized Defensive Task Lists:\n- Remediations are sorted automatically based on predicted severity scores, tackling high-impact risks first."
    ]
    create_unified_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.6), "Targeted Defense & Remediations", rem_bullets, COLOR_CYAN, space_after=8)

    # -------------------------------------------------------------
    # SLIDE 8: Security Architecture & Best Practices
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "Defensive Engineering & Integrity")
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Architectural principles implemented to secure the vulnerability management platform itself:"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_TEXT_SECONDARY

    sec_cards = [
        {"title": "Solid Input Sanitization", "bullets": ["AGGRESSIVE sanitizations using target regex patterns.", "Filters remote commands, avoiding parameter injections."]},
        {"title": "SQLi & CSRF Shields", "bullets": ["SQLAlchemy ORM blocks SQL injections.", "CSRFProtect locks forms, blocking cross-site requests."]},
        {"title": "Bcrypt Cryptography", "bullets": ["User credentials salted and hashed using bcrypt libraries.", "Secured session scopes with cryptographically signed tokens."]},
        {"title": "Isolated Executions", "bullets": ["Containerized setups wrapping active scan environments.", "Avoids local command line overrides during target probes."]}
    ]
    
    for i, c in enumerate(sec_cards):
        left_pos = Inches(0.8 + i * 2.95)
        create_unified_card(slide, left_pos, Inches(2.2), Inches(2.8), Inches(4.4), c["title"], c["bullets"], COLOR_CYAN, space_after=8)

    # -------------------------------------------------------------
    # SLIDE 9: Platform Deployment & Stack
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "Technical Stack & Deployments")
    
    # Left: Stack Breakdown
    bullets_stack = [
        "Backend Infrastructure:\n- Python 3.13 & Flask App Factory blueprints.\n- SQLite databases utilizing SQLAlchemy ORM schema maps.",
        "Frontend Presentation Console:\n- HTML5, CSS3 Glassmorphism variables, Bootstrap 5.\n- Chart.js metrics engines and Typewriter JS animations.",
        "AI/ML Diagnostic Pipeline:\n- Scikit-Learn (Random Forest Classifiers), Pandas, NumPy.",
        "Reporting Systems:\n- ReportLab programmatic drawing compilers."
    ]
    create_unified_card(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6), "The Engineering Architecture", bullets_stack, COLOR_TEAL, space_after=4)
    
    # Right: Deployment
    bullets_deploy = [
        "Sandboxed Containerization & Setup:\n- Custom Dockerfile wrapping Python runtimes, Scikit-learn pipelines, and Nmap OS binaries directly inside a safe sandbox.",
        "Two-Step Quick Start:\n- 1. Build container: docker build -t securescan:latest .\n- 2. Deploy sandbox: docker run -d -p 5000:5000 --name securescan securescan:latest",
        "Resilient Local Backups:\n- Runs seamlessly in network simulation mode on restricted corporate systems or offline development spaces."
    ]
    create_unified_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.6), "Sandboxed Containerization & Setup", bullets_deploy, COLOR_CYAN, space_after=4)

    # -------------------------------------------------------------
    # SLIDE 10: Future Vision & Roadmap
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    add_header(slide, "Future Roadmap & Vision")
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Scaling SecureScan from a localized scanner into an enterprise security mesh:"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_TEXT_SECONDARY

    # Roadmap Steps (3 Horizontal Cards)
    roadmap = [
        {"step": "PHASE 01", "name": "Distributed Agents", "desc": "Deploy low-footprint scanning daemons across network segments. Centralize all diagnostics into the core Flask engine."},
        {"step": "PHASE 02", "name": "Live CVE Feeds", "desc": "Establish real-time API sync tunnels with the National Vulnerability Database (NVD) to fetch live CVSS scores daily."},
        {"step": "PHASE 03", "name": "Generative LLM Patching", "desc": "Integrate a sandboxed Local LLM to compile concrete source code patches for scanned application vulnerabilities."}
    ]
    
    for i, r in enumerate(roadmap):
        left_pos = Inches(0.8 + i * 4.0)
        create_roadmap_card(slide, left_pos, Inches(2.2), Inches(3.7), Inches(4.4), r["step"], r["name"], r["desc"])

    # Save Presentation
    filename = "SecureScan_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as '{filename}'!")

if __name__ == "__main__":
    main()
